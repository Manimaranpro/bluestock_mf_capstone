import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Simplified Presentation Layout helper
class PowerPointBuilder:
    def __init__(self, prs, DARK_BLUE, CHARCOAL):
        self.prs = prs
        self.DARK_BLUE = DARK_BLUE
        self.CHARCOAL = CHARCOAL
        self.slides = prs.slides
        self.slide_layouts = prs.slide_layouts

    def build_slide(self, data, idx):
        if idx == 0:
            slide = self.slides.add_slide(self.slide_layouts[5]) # Blank layout
            
            # Title Box
            txBox = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.3), Inches(4.0))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = data["title"]
            p.font.size = Pt(40)
            p.font.bold = True
            p.font.color.rgb = self.DARK_BLUE
            p.space_after = Pt(20)
            
            p2 = tf.add_paragraph()
            p2.text = data["subtitle"]
            p2.font.size = Pt(18)
            p2.font.color.rgb = self.CHARCOAL
        else:
            slide = self.slides.add_slide(self.slide_layouts[5]) # Blank layout
            
            # Title
            txBox_title = slide.shapes.add_textbox(Inches(1.0), Inches(0.6), Inches(11.3), Inches(1.0))
            tf_title = txBox_title.text_frame
            tf_title.word_wrap = True
            p_title = tf_title.paragraphs[0]
            p_title.text = data["title"]
            p_title.font.size = Pt(28)
            p_title.font.bold = True
            p_title.font.color.rgb = self.DARK_BLUE
            
            # Content Bullets
            txBox_content = slide.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.3), Inches(5.0))
            tf_content = txBox_content.text_frame
            tf_content.word_wrap = True
            for i, bullet in enumerate(data["bullets"]):
                p = tf_content.paragraphs[0] if i == 0 else tf_content.add_paragraph()
                p.text = f"•   {bullet}"
                p.font.size = Pt(16)
                p.font.color.rgb = self.CHARCOAL
                p.space_after = Pt(18)

if __name__ == "__main__":
    # Create the builder instance
    prs_obj = Presentation()
    prs_obj.slide_width = Inches(13.33)
    prs_obj.slide_height = Inches(7.5)
    
    DARK_BLUE = RGBColor(0, 51, 102)
    CHARCOAL = RGBColor(51, 51, 51)
    
    builder = PowerPointBuilder(prs_obj, DARK_BLUE, CHARCOAL)
    
    # We populate the structural slide array
    slides_list = [
        {
            "title": "Bluestock Mutual Fund Capstone",
            "subtitle": "Quantitative Portfolio Performance, Risk Modeling, & Investor Continuity Analytics\n\nPrepared by: Data Analytics Intern\nSubmission Date: July 1, 2026\nProject Manager: Yash Kale (yashkale@bluestock.in)"
        },
        {
            "title": "Technical ETL & Database Architecture",
            "bullets": [
                "Modular ETL pipeline developed using dynamic relative paths for complete environment portability.",
                "Non-pricing weekend and holiday gaps resolved by reindexing daily date ranges and applying forward-fill (.ffill()) algorithms.",
                "Cleaned tables structured in a star schema and loaded to SQLite (bluestock_mf.db) to optimize analytic query execution."
            ]
        },
        {
            "title": "Fund Performance Scorecard Leadership",
            "bullets": [
                "Constructed a 100-point multi-factor scorecard combining 3Yr CAGR, Sharpe, Alpha, Expense, and Drawdown rankings.",
                "Rank #1: Mirae Asset Tax Saver Fund - Regular - Growth (Score: 87.94 | Sharpe: 1.23 | Sortino: 2.15 | Alpha: 28.08%).",
                "Rank #2: DSP Small Cap Fund - Regular - Growth (Score: 87.44 | Sharpe: 0.95 | Alpha: 29.58% | Drawdown: -31.2%).",
                "Rank #3: ICICI Pru Midcap Fund - Regular - Growth (Score: 82.38 | Sharpe: 1.18 | Sortino: 2.03 | Alpha: 29.51%)."
            ]
        },
        {
            "title": "Quantitative Risk Modeling (VaR vs. CVaR)",
            "bullets": [
                "Modeled daily historical downside risk over 1,149 active trading days using 95% confidence intervals.",
                "Extreme Tail Loss: SBI Small Cap Fund (AMFI: 119599) has a daily VaR of -2.69% and a Conditional VaR (CVaR) of -3.24%.",
                "Capital Preservation: ICICI Pru Liquid Fund (AMFI: 120507) has a daily VaR of -0.02% and a Conditional VaR of -0.04%."
            ]
        },
        {
            "title": "Diversification & Sector Concentration",
            "bullets": [
                "Evaluated asset concentration across equity portfolios using the Herfindahl-Hirschman Index (HHI).",
                "Axis Bluechip Fund (AMFI: 119092): Highest HHI sector concentration of 0.2968, indicating sensitivity to financial/tech corrections.",
                "UTI Mid Cap Fund (AMFI: 102886): Lowest HHI sector concentration of 0.1240, indicating highly diversified industry holdings."
            ]
        },
        {
            "title": "Long-Term NAV Projection (B3 Monte Carlo)",
            "bullets": [
                "Executed a 5-year stochastic simulation (1,260 future trading days) with 1,000 runs using historical volatility.",
                "Projected Mirae Asset Tax Saver Fund expected NAV paths, bounding extreme best and worst-case outcomes.",
                "Visual chart saved to reports/monte_carlo_projection.png to help advisors set objective long-term growth expectations."
            ]
        },
        {
            "title": "Investor Behavioral & Retention Analysis",
            "bullets": [
                "Audit of loyal accounts (6+ consecutive SIPs) flagged payment intervals exceeding 35 days as high-retention risks.",
                "At-Risk Accounts: INV000028 (avg gap: 93.6 days) and INV000004 (avg gap: 85.4 days) show extensive mandate friction.",
                "Disciplined Accounts: INV000334 (gap of 25.5 days) and INV000750 (gap of 26.5 days) represent stable recurring inflows."
            ]
        },
        {
            "title": "Interactive Analytics Dashboard Overview",
            "bullets": [
                "Developed a 4-page Power BI dashboard detailing Industry Overview, Performance, Investor Analytics, and SIP Trends.",
                "Configured multi-select interactive slicers (State, Fund House, Age Group, Category) for real-time portfolio filtering.",
                "Verified industry AUM baseline scales of ₹81L Cr, monthly SIP inflows of ₹31K Cr, and a folio base of 26.12 Cr."
            ]
        },
        {
            "title": "Strategic Recommendations",
            "bullets": [
                "Establish automated UI warnings flag limiters for high-concentration equity funds (sector HHI > 0.20).",
                "Deploy proactive notification workflows targeting at-risk accounts whose transactional gap has exceeded 35 days.",
                "Incorporate risk-adjusted metrics (Sharpe, Sortino, Alpha) directly into the core scheme recommendation logic."
            ]
        }
    ]

    for i, slide_data in enumerate(slides_list):
        builder.build_slide(slide_data, i)

    # Resolve output directory paths
    from pathlib import Path
    SCRIPT_DIR = Path(__file__).resolve().parent
    PROJECT_ROOT = SCRIPT_DIR.parent
    reports_dir = PROJECT_ROOT / "reports"
    os.makedirs(reports_dir, exist_ok=True)
    
    output_path = reports_dir / "Presentation.pptx"
    prs_obj.save(output_path)
    print(f"Presentation slide deck successfully generated at: {output_path}")